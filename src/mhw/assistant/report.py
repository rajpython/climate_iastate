"""PowerPoint export — turn a chat's charts/tables/findings into a downloadable ``.pptx`` deck.

Server-side and frontend-agnostic: returns a token; the API serves the file by token. Charts are
Plotly figure dicts (from :mod:`charts`) rasterised to PNG via kaleido — **pinned kaleido 0.2.1**
because 1.x needs a system Chrome the slim image lacks (0.2.x bundles its own).

Deck layout: a title slide, an auto exec-summary (when there are ≥2 content slides), then one slide
per item — a blue heading, a large chart OR a native table on the left, interpretation bullets on the
right, and a branded footer. The chart's own title is dropped (the slide heading is the title).
"""
from __future__ import annotations

import html
import os
import tempfile
import uuid
from pathlib import Path
from typing import Any


def _txt(v) -> str:
    """Plain text for pptx: decode HTML entities the model sometimes emits (e.g. '&amp;' → '&')."""
    return html.unescape("" if v is None else str(v))

REPORTS_DIR = Path(
    os.getenv("ASSISTANT_REPORTS_DIR", str(Path(tempfile.gettempdir()) / "mhw_assistant_reports"))
)

_HEADING_RGB = (0x16, 0x40, 0x7A)   # NOAA-style blue, matches the dashboard
_BODY_RGB = (0x33, 0x41, 0x4F)      # dark slate
_FOOTER_RGB = (0x9A, 0xA4, 0xB0)
_HDR_FILL = (0x16, 0x40, 0x7A)
_FOOTER = "Alaska Marine Ecosystems Dashboard · marine.iastate.ai"


def report_path(token: str) -> Path | None:
    """Resolve a download token to a file path (None if unknown/absent). Token is a bare filename."""
    if not token or "/" in token or "\\" in token or ".." in token:
        return None
    p = REPORTS_DIR / token
    return p if p.exists() else None


def _fig_png(fig_dict: dict) -> bytes:
    import plotly.graph_objects as go
    import plotly.io as pio

    fig = go.Figure(fig_dict)
    # The slide heading is the title — drop the chart's own title and reclaim the top space.
    fig.update_layout(title_text="", margin={"t": 28})
    return pio.to_image(fig, format="png", width=1120, height=680, scale=2)


def build_report(
    title: str,
    slides: list[dict[str, Any]],
    subtitle: str = "",
    out_dir: Path | None = None,
) -> dict:
    """Build a deck and return ``{"token", "filename"}``.

    Each slide is ``{"heading", "bullets": [str], "chart": <plotly dict>|None,
    "table": {"columns": [...], "rows": [[...]]}|None}``.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    out_dir = Path(out_dir) if out_dir else REPORTS_DIR
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _footer(slide, idx: int) -> None:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
        p = box.text_frame.paragraphs[0]
        p.text = f"{_FOOTER}     ·     {idx}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(*_FOOTER_RGB)

    def _heading(slide, text: str) -> None:
        head = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        htf = head.text_frame
        htf.word_wrap = True
        htf.text = _txt(text)
        hp = htf.paragraphs[0]
        hp.font.size = Pt(26)
        hp.font.bold = True
        hp.font.color.rgb = RGBColor(*_HEADING_RGB)

    def _bullets(box, items, size_pt):
        btf = box.text_frame
        btf.word_wrap = True
        for i, b in enumerate(items):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.text = f"•  {_txt(b)}"
            p.font.size = Pt(size_pt)
            p.font.color.rgb = RGBColor(*_BODY_RGB)
            p.space_after = Pt(9)
            p.line_spacing = 1.08

    def _table(slide, columns, rows, left, top, width):
        rows = rows[:14]                       # keep decks readable
        gt = slide.shapes.add_table(len(rows) + 1, len(columns), left, top, width,
                                    Inches(0.35 * (len(rows) + 1))).table
        for c, col in enumerate(columns):
            cell = gt.cell(0, c)
            cell.text = _txt(col)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*_HDR_FILL)
        for r, row in enumerate(rows, start=1):
            for c in range(len(columns)):
                cell = gt.cell(r, c)
                cell.text = "" if c >= len(row) or row[c] is None else _txt(row[c])
                cell.text_frame.paragraphs[0].font.size = Pt(11)

    # --- Title slide ---
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = _txt(title) or "Alaska Marine Ecosystems — Data Report"
    if s0.placeholders and len(s0.placeholders) > 1:
        s0.placeholders[1].text = _txt(subtitle) or "Generated from marine.iastate.ai"

    # --- Auto exec-summary (only when it adds value) ---
    if len([s for s in slides if isinstance(s, dict)]) >= 2:
        s = prs.slides.add_slide(blank)
        _heading(s, "Summary")
        box = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.2))
        _bullets(box, [str(sl.get("heading", "")) for sl in slides if isinstance(sl, dict)], 18)
        _footer(s, 1)

    # --- Content slides ---
    for idx, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        s = prs.slides.add_slide(blank)
        _heading(s, slide.get("heading", ""))
        chart, table, bullets = slide.get("chart"), slide.get("table"), slide.get("bullets") or []

        if chart:
            img_path = out_dir / f"_img_{uuid.uuid4().hex}.png"
            img_path.write_bytes(_fig_png(chart))
            s.shapes.add_picture(str(img_path), Inches(0.45), Inches(1.35), width=Inches(8.5))
            img_path.unlink(missing_ok=True)
            _bullets(s.shapes.add_textbox(Inches(9.25), Inches(1.45), Inches(3.65), Inches(5.4)),
                     bullets, 14)
        elif table and table.get("rows"):
            _table(s, table.get("columns", []), table["rows"], Inches(0.5), Inches(1.4), Inches(8.3))
            _bullets(s.shapes.add_textbox(Inches(9.25), Inches(1.45), Inches(3.65), Inches(5.4)),
                     bullets, 14)
        else:
            _bullets(s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.3)),
                     bullets, 18)
        _footer(s, idx + 1)

    token = f"{uuid.uuid4().hex}.pptx"
    prs.save(str(out_dir / token))
    safe_name = (title or "report").strip().replace(" ", "_")[:60] or "report"
    return {"token": token, "filename": f"{safe_name}.pptx"}
